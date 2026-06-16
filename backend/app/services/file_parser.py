from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass
from io import BytesIO

from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

from ..exceptions import TextExtractionError


@dataclass(slots=True)
class ExtractResult:
	text: str
	filename: str
	file_type: str
	page_count: int
	char_count: int
	word_count: int
	truncated: bool
	warning: str | None


class FileParser:
	def __init__(self, *, max_text_length: int):
		self._max_text_length = max_text_length

	def parse(self, *, file_bytes: bytes, filename: str) -> ExtractResult:
		file_type = self._detect_file_type(filename)
		if file_type == "pdf":
			raw_text, page_count = self._parse_pdf(file_bytes)
		else:
			raw_text, page_count = self._parse_pptx(file_bytes)

		cleaned_text, truncated = self._clean_text(raw_text)
		char_count = len(cleaned_text)
		word_count = len(cleaned_text.split()) if cleaned_text else 0

		warning = None
		if len(cleaned_text.strip()) < 50:
			warning = (
				"Very little text was extracted. The file may contain mostly images or scanned content."
			)

		return ExtractResult(
			text=cleaned_text,
			filename=filename,
			file_type=file_type,
			page_count=page_count,
			char_count=char_count,
			word_count=word_count,
			truncated=truncated,
			warning=warning,
		)

	@staticmethod
	def _detect_file_type(filename: str) -> str:
		lowered = (filename or "").lower()
		if lowered.endswith(".pdf"):
			return "pdf"
		if lowered.endswith(".pptx"):
			return "pptx"
		raise TextExtractionError("Only PDF and PPTX files are supported")

	@staticmethod
	def _parse_pdf(file_bytes: bytes) -> tuple[str, int]:
		try:
			reader = PdfReader(BytesIO(file_bytes))
		except PdfReadError as exc:
			raise TextExtractionError("The PDF file appears to be corrupted or invalid.") from exc
		except Exception as exc:  # noqa: BLE001
			raise TextExtractionError(
				"An error occurred while processing the PDF. Please try a different file."
			) from exc

		if not reader.pages:
			raise TextExtractionError("The PDF file contains no pages.")

		if getattr(reader, "is_encrypted", False):
			try:
				decrypt_result = reader.decrypt("")
			except Exception as exc:  # noqa: BLE001
				raise TextExtractionError(
					"This PDF is password-protected. Please upload an unprotected version."
				) from exc
			if decrypt_result == 0:
				raise TextExtractionError(
					"This PDF is password-protected. Please upload an unprotected version."
				)

		chunks: list[str] = []
		for idx, page in enumerate(reader.pages, start=1):
			try:
				text = page.extract_text() or ""
			except Exception:
				text = ""
			text = text.strip()
			if not text:
				continue
			chunks.append(f"--- Page {idx} ---\n{text}")

		return "\n\n".join(chunks), len(reader.pages)

	@staticmethod
	def _parse_pptx(file_bytes: bytes) -> tuple[str, int]:
		try:
			pres = Presentation(BytesIO(file_bytes))
		except PackageNotFoundError as exc:
			raise TextExtractionError("The PPTX file appears to be corrupted or invalid.") from exc
		except ValueError as exc:
			raise TextExtractionError("Could not parse the presentation file.") from exc
		except Exception as exc:  # noqa: BLE001
			raise TextExtractionError(
				"An error occurred while processing the PPTX. Please try a different file."
			) from exc

		slide_count = len(pres.slides)
		chunks: list[str] = []

		def iter_shape_text(shape) -> list[str]:
			parts: list[str] = []
			try:
				if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
					for sub in shape.shapes:
						parts.extend(iter_shape_text(sub))
					return parts
			except Exception:
				pass

			if getattr(shape, "has_table", False):
				try:
					for row in shape.table.rows:
						for cell in row.cells:
							cell_text = (cell.text or "").strip()
							if cell_text:
								parts.append(cell_text)
				except Exception:
					pass

			if getattr(shape, "has_text_frame", False):
				try:
					for paragraph in shape.text_frame.paragraphs:
						para_text = (paragraph.text or "").strip()
						if para_text:
							parts.append(para_text)
				except Exception:
					pass

			return parts

		for idx, slide in enumerate(pres.slides, start=1):
			slide_parts: list[str] = []

			try:
				title_shape = slide.shapes.title
				title_text = (title_shape.text or "").strip() if title_shape else ""
				if title_text:
					slide_parts.append(title_text)
			except Exception:
				pass

			for shape in slide.shapes:
				slide_parts.extend(iter_shape_text(shape))

			slide_text = "\n".join([p for p in slide_parts if p]).strip()
			if not slide_text:
				continue
			chunks.append(f"--- Slide {idx} ---\n{slide_text}")

		return "\n\n".join(chunks), slide_count

	def _clean_text(self, raw_text: str) -> tuple[str, bool]:
		text = raw_text or ""
		if not text:
			return "", False

		text = text.replace("\r\n", "\n").replace("\r", "\n")
		text = unicodedata.normalize("NFKC", text)

		# Remove null bytes / control chars (except \n and \t)
		text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)

		lines = [ln.strip() for ln in text.split("\n")]

		# Remove pure page-number artifacts
		page_num_patterns = [
			re.compile(r"^\d+$"),
			re.compile(r"^-+\s*\d+\s*-+$"),
			re.compile(r"^page\s+\d+$", re.IGNORECASE),
		]

		filtered: list[str] = []
		for ln in lines:
			if not ln:
				filtered.append("")
				continue
			if ln.startswith("--- Page") or ln.startswith("--- Slide"):
				filtered.append(ln)
				continue
			if any(pat.match(ln) for pat in page_num_patterns):
				continue
			# Collapse multiple spaces/tabs within the line
			ln = re.sub(r"[ \t]{2,}", " ", ln)
			filtered.append(ln)

		# Remove repeated short header/footer-like lines (basic heuristic)
		counts: dict[str, int] = {}
		for ln in filtered:
			if not ln:
				continue
			if ln.startswith("--- Page") or ln.startswith("--- Slide"):
				continue
			if len(ln) <= 30:
				counts[ln] = counts.get(ln, 0) + 1

		repeated = {ln for ln, cnt in counts.items() if cnt >= 3}
		filtered = ["" if ln in repeated else ln for ln in filtered]

		text = "\n".join(filtered)
		text = re.sub(r"\n{3,}", "\n\n", text).strip()

		truncated = False
		if len(text) > self._max_text_length:
			cutoff = text.rfind("\n\n", 0, self._max_text_length)
			if cutoff <= 0:
				cutoff = self._max_text_length
			text = text[:cutoff].rstrip()
			text += (
				"\n\n[Note: Text was truncated due to length. Only the first ~"
				f"{self._max_text_length} characters were included.]"
			)
			truncated = True

		return text, truncated
